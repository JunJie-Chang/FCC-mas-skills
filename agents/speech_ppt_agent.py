"""
agents/speech_ppt_agent.py — Speech PPT generator using LangGraph.

Flow:
    parse_script → confirm_slides → generate_images → build_ppt

Input:  CY's verbal transcript describing slide content (raw_transcript /
        task_instruction).  Audio is transcribed upstream (STT in main.py)
        or by the standalone CLI (--audio).

Slides are classified as:
  - structured:   standard layout (title + 5 bullets + right-side image)
  - unstructured: complex layouts (charts, tables, org-charts, etc.);
                  CY describes these verbally — notes are echoed back to
                  the operator for manual handling; no PPT slide is built.

Output: .pptx file (structured slides only) + notes printed to console
        and appended to the .log sidecar.

Filename: YYYY.MM.DD_SpeechTopic_InternName.pptx

Reference layout: Works/2026.04.09_台中智慧製造演講.pptx
  Slide size:   10.0 × 7.5 inches
  Title:        pos=(0.51, 0.40) size=(8.98, 0.93) 28pt
  Bullet area:  pos=(0.51, 1.49) size=(4.70, 4.31) 24pt bold
  Image:        pos=(5.45, 1.49) size=(4.20, 4.31)
  Page num:     pos=(4.23, 6.95) size=(1.54, 0.40) 11pt centred
"""
import json
import os
import sys
import tempfile
from datetime import date
from pathlib import Path
from typing import TypedDict, Union

from anthropic import Anthropic
from dotenv import load_dotenv
from langgraph.graph import END, StateGraph

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

import config
from utils.file_naming import general
from utils.logger import AgentLogger

load_dotenv(override=True)


# ── Layout constants (inches → EMU: 1in = 914400) ────────────────────────────

_IN = 914400

SLIDE_W  = int(10.0 * _IN)
SLIDE_H  = int(7.5  * _IN)

TITLE_L  = int(0.51 * _IN)
TITLE_T  = int(0.40 * _IN)
TITLE_W  = int(8.98 * _IN)
TITLE_H  = int(0.93 * _IN)
TITLE_PT = 28

BULLET_L = int(0.51 * _IN)
BULLET_T = int(1.49 * _IN)
BULLET_W = int(4.70 * _IN)
BULLET_H = int(4.31 * _IN)
BULLET_PT = 24

IMAGE_L  = int(5.45 * _IN)
IMAGE_T  = int(1.49 * _IN)
IMAGE_W  = int(4.20 * _IN)
IMAGE_H  = int(4.31 * _IN)

PAGENUM_L = int(4.23 * _IN)
PAGENUM_T = int(6.95 * _IN)
PAGENUM_W = int(1.54 * _IN)
PAGENUM_H = int(0.40 * _IN)

# Chrome template (0-slide PPTX in repo — preserves master/layout chrome)
_CHROME_TEMPLATE = Path(__file__).resolve().parents[1] / "assets" / "ppt_chrome_template.pptx"

# Source for one-time chrome template creation (OneDrive, may be absent on other machines)
_PPT_REFERENCE = Path(
    "/Users/junjie/Library/CloudStorage/OneDrive-個人(2)/Internship/Works/"
    "2026.04.09_台中智慧製造演講.pptx"
)


# ── Style helpers ─────────────────────────────────────────────────────────────

def _ensure_chrome_template() -> bool:
    """
    Create assets/ppt_chrome_template.pptx from the reference PPTX if it doesn't
    exist yet.  Returns True if the template is ready, False if unavailable.
    """
    if _CHROME_TEMPLATE.exists():
        return True
    if not _PPT_REFERENCE.exists():
        return False
    try:
        from pptx import Presentation as _Prs
        prs = _Prs(str(_PPT_REFERENCE))
        r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        sldIdLst = prs.slides._sldIdLst
        for sldId in list(sldIdLst):
            r_id = sldId.get(f"{{{r_ns}}}id")
            if hasattr(prs.part, "_rels") and hasattr(prs.part._rels, "_rels"):
                prs.part._rels._rels.pop(r_id, None)
            sldIdLst.remove(sldId)
        _CHROME_TEMPLATE.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(_CHROME_TEMPLATE))
        print(f"  [setup] Chrome template created: {_CHROME_TEMPLATE}")
        return True
    except Exception as e:
        print(f"  [WARN] Could not create chrome template: {e}")
        return False


def _get_object_layout(prs):
    """Return the 'OBJECT' slide layout from the presentation's first master."""
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name == "OBJECT":
            return layout
    return prs.slide_layouts[6]  # fallback: blank


def _set_bullet_para_props(para) -> None:
    """
    Apply paragraph formatting matching the reference PPTX:
    150% line spacing, 10pt space-before, hanging indent.
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    pPr = para._p.get_or_add_pPr()
    pPr.set("marL", "228600")
    pPr.set("indent", "-228600")
    pPr.set("eaLnBrk", "0")
    pPr.set("fontAlgn", "base")
    pPr.set("hangingPunct", "0")

    # Remove any existing lnSpc / spcBef to avoid duplicates
    for tag in ("a:lnSpc", "a:spcBef"):
        old = pPr.find(qn(tag))
        if old is not None:
            pPr.remove(old)

    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    etree.SubElement(lnSpc, qn("a:spcPct")).set("val", "150000")

    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(spcBef, qn("a:spcPts")).set("val", "1000")

    # White bullet character (•)
    buClr = etree.SubElement(pPr, qn("a:buClr"))
    etree.SubElement(buClr, qn("a:srgbClr")).set("val", "FFFFFF")
    etree.SubElement(pPr, qn("a:buSzPct")).set("val", "100000")
    etree.SubElement(pPr, qn("a:buChar")).set("char", "•")


# ── State ─────────────────────────────────────────────────────────────────────

class SpeechPPTState(TypedDict):
    speech_topic: str                  # for filename (derived or explicit)
    raw_transcript: str                # CY's verbal description
    intern_name: Union[str, list[str]]
    task_date: str
    subdir: str
    generate_images: bool
    slides_plan: list[dict]            # [{type, title, bullets?, notes?}]
    confirmed: bool
    output_path: str
    log_path: str


# ── Helpers ───────────────────────────────────────────────────────────────────

def _get_anthropic() -> Anthropic:
    key = os.getenv("ANTHROPIC_API_KEY")
    if not key:
        raise EnvironmentError("ANTHROPIC_API_KEY not set in .env")
    return Anthropic(api_key=key)


def _get_openai():
    from openai import OpenAI
    key = os.getenv("OPENAI_API_KEY")
    if not key:
        raise EnvironmentError("OPENAI_API_KEY not set in .env")
    return OpenAI(api_key=key)


def _parse_json(text: str):
    import re
    original = text
    text = text.strip()
    text = re.sub(r"^```(?:json)?\s*", "", text)
    text = re.sub(r"\s*```$", "", text)
    text = text.strip()
    match = re.search(r"(\{[\s\S]*\}|\[[\s\S]*\])", text)
    if match:
        text = match.group(1)
    try:
        return json.loads(text)
    except json.JSONDecodeError as e:
        print(f"\n[DEBUG] JSON parse failed: {e}")
        print(f"[DEBUG] Raw response:\n{original[:2000]}")
        raise


def _intern_str(intern_name) -> str:
    if isinstance(intern_name, list):
        return ", ".join(intern_name)
    return intern_name or config.DEFAULT_INTERN_NAME


# ── Node 1: parse_script ──────────────────────────────────────────────────────

def parse_script(state: SpeechPPTState) -> dict:
    """
    Parse CY's transcript into a slides_plan list.
    Each entry is either:
      {"type": "structured", "title": str, "bullets": [5 strs]}
      {"type": "unstructured", "title": str, "notes": str}
    Also derives speech_topic if not already set.
    """
    client = _get_anthropic()
    transcript = state["raw_transcript"]
    topic = state.get("speech_topic", "") or ""

    prompt = f"""你是一位演講簡報規劃助理。根據以下演講內容描述（錄音轉文字），整理出每張投影片的規劃。

【分類規則】
- 「structured」：標準版面頁（標題 + 5 條重點條列 + 右側配圖），系統自動生成
- 「unstructured」：非標準版面頁（圖表、比較表、組織圖、流程圖、引用頁等），需操作者手動處理

【structured 頁面要求】
- title：投影片標題，≤ 15 字
- bullets：5 條重點，每條 ≤ 10 個中文字（英文單字算 1 個字）

【unstructured 頁面要求】
- title：投影片標題，≤ 15 字
- notes：完整保留 CY 的原始說明，不得改寫或刪減

【判斷為非結構化的線索】
- CY 提到「配合紙本」、「這張比較複雜」、「放一個圖」、「比較表」、「圓圈」、「箭頭」
- 描述涉及多欄位、空間佈局、圖表關係
- 引用他人的話語或特定數據來源頁面

演講主題：{topic or "（請自行判斷）"}

演講描述：
{transcript}

請回傳純 JSON（不要加說明文字）：
{{
  "speech_title": "演講主題名稱（≤15字；若上方已提供可省略此欄）",
  "slides": [
    {{"type": "structured", "title": "投影片標題", "bullets": ["條列1", "條列2", "條列3", "條列4", "條列5"]}},
    {{"type": "unstructured", "title": "投影片標題", "notes": "CY 原始說明完整內容..."}}
  ]
}}"""

    msg = client.messages.create(
        model=config.LLM_MAIN,
        max_tokens=4096,
        messages=[{"role": "user", "content": prompt}],
    )
    from utils.cost_tracker import tracker
    tracker.record_claude(config.LLM_MAIN, msg.usage.input_tokens, msg.usage.output_tokens)

    raw = next((b.text for b in msg.content if hasattr(b, "text")), "")
    data = _parse_json(raw)

    slides_plan = data.get("slides", [])
    derived_title = data.get("speech_title") or topic or "演講"
    final_topic = topic or derived_title

    return {
        "slides_plan": slides_plan,
        "speech_topic": final_topic,
    }


# ── Node 2: confirm_slides ────────────────────────────────────────────────────

def confirm_slides(state: SpeechPPTState) -> dict:
    """
    Show slide plan and wait for CLI confirmation before generating images.
    User must confirm before DALL-E costs are incurred.
    """
    slides = state["slides_plan"]
    topic = state["speech_topic"]
    do_images = state.get("generate_images", True)

    structured = [s for s in slides if s.get("type") == "structured"]
    unstructured = [s for s in slides if s.get("type") != "structured"]

    print(f"\n── Speech PPT 計劃：{topic} {'─' * 28}")
    for i, slide in enumerate(slides, 1):
        tag = "結構化" if slide.get("type") == "structured" else "非結構"
        print(f"\n  [{i}]  {tag}  {slide['title']}")
        if slide.get("type") == "structured":
            for b in slide.get("bullets", []):
                print(f"         • {b}")
        else:
            notes = slide.get("notes", "")
            lines = notes.strip().splitlines()
            for line in lines[:3]:
                print(f"         📝 {line}")
            if len(lines) > 3:
                print(f"         📝 ...")

    img_note = f"，將生成 {len(structured)} 張 DALL-E 圖片" if do_images and structured else ""
    print(f"\n  共 {len(slides)} 張：{len(structured)} 結構化（生成 PPT{img_note}）"
          f"、{len(unstructured)} 非結構（僅回傳 notes）")
    print("─" * 50)

    while True:
        ans = input("確認執行？(y = 確認 / n = 取消)：").strip().lower()
        if ans in ("y", "yes", ""):
            return {"confirmed": True}
        if ans in ("n", "no"):
            print("已取消。")
            return {"confirmed": False}
        print("請輸入 y 或 n。")


def _route_after_confirm(state: SpeechPPTState) -> str:
    return "generate_images" if state.get("confirmed") else END


# ── Node 3: generate_images ───────────────────────────────────────────────────

def generate_images(state: SpeechPPTState) -> dict:
    """Generate DALL-E 3 images for structured slides; skip unstructured."""
    if not state.get("generate_images", True):
        return {}

    client = _get_openai()
    slides = state["slides_plan"]
    MAX_RETRIES = 3

    for i, slide in enumerate(slides):
        if slide.get("type") != "structured":
            continue
        title = slide["title"]
        print(f"  DALL-E [{i + 1}]: {title}")
        for attempt in range(1, MAX_RETRIES + 1):
            try:
                prompt = (
                    f"Photorealistic photograph or realistic illustration representing: '{title}'. "
                    "Professional corporate context. Absolutely no text, words, letters, "
                    "numbers, signs, labels, or watermarks anywhere in the image. "
                    "Suitable as a right-side panel image on a dark navy PowerPoint slide."
                )
                resp = client.images.generate(
                    model="dall-e-3",
                    prompt=prompt,
                    size="1024x1024",
                    quality="standard",
                    n=1,
                )
                image_url = resp.data[0].url
                from utils.cost_tracker import tracker
                tracker.record_dalle(1)

                import ssl
                import urllib.request
                ctx = ssl._create_unverified_context()
                tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                with urllib.request.urlopen(image_url, context=ctx) as resp:
                    tmp.write(resp.read())
                tmp.close()
                slide["image_path"] = tmp.name
                break
            except Exception as e:
                if attempt < MAX_RETRIES:
                    print(f"  [WARN] retry {attempt}/{MAX_RETRIES} for '{title}': {e}")
                else:
                    print(f"  [WARN] image generation failed after {MAX_RETRIES} attempts for '{title}'")
                    slide["image_path"] = None

    return {"slides_plan": slides}


# ── Node 4: build_ppt ─────────────────────────────────────────────────────────

def build_ppt(state: SpeechPPTState) -> dict:
    """
    Build PPTX for structured slides.
    Print notes for unstructured slides and append to .log sidecar.
    """
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor

    import shutil

    slides = state["slides_plan"]
    structured   = [s for s in slides if s.get("type") == "structured"]
    unstructured = [s for s in slides if s.get("type") != "structured"]

    # Open chrome template (preserves master background, line, logo, page-num placeholder)
    # Fall back to a blank presentation if the template is unavailable.
    if _ensure_chrome_template():
        prs = Presentation(str(_CHROME_TEMPLATE))
        object_layout = _get_object_layout(prs)
    else:
        print(f"  [WARN] Chrome template not found: {_CHROME_TEMPLATE}")
        prs = Presentation()
        prs.slide_width  = Emu(SLIDE_W)
        prs.slide_height = Emu(SLIDE_H)
        object_layout = prs.slide_layouts[6]

    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    for slide_num, slide in enumerate(structured, start=1):
        sl = prs.slides.add_slide(object_layout)
        # Chrome elements (background, blue line, logo, sldNum placeholder) are
        # inherited from the OBJECT layout — no manual additions needed.

        # Title — textbox overlaid on the layout's title placeholder area
        title_box = sl.shapes.add_textbox(
            Emu(TITLE_L), Emu(TITLE_T), Emu(TITLE_W), Emu(TITLE_H)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = slide["title"]
        run.font.size = Pt(TITLE_PT)
        run.font.bold = False
        run.font.color.rgb = WHITE

        # Bullets — left-half textbox (right half reserved for image)
        bullet_box = sl.shapes.add_textbox(
            Emu(BULLET_L), Emu(BULLET_T), Emu(BULLET_W), Emu(BULLET_H)
        )
        btf = bullet_box.text_frame
        btf.word_wrap = True
        for bi, bullet in enumerate(slide.get("bullets", [])[:5]):
            bp = btf.paragraphs[0] if bi == 0 else btf.add_paragraph()
            _set_bullet_para_props(bp)
            brun = bp.add_run()
            brun.text = bullet
            brun.font.size = Pt(BULLET_PT)
            brun.font.bold = True
            brun.font.color.rgb = WHITE

        # Content image (right half)
        img_path = slide.get("image_path")
        if img_path and Path(img_path).exists():
            sl.shapes.add_picture(
                img_path,
                Emu(IMAGE_L), Emu(IMAGE_T), Emu(IMAGE_W), Emu(IMAGE_H),
            )

    # Save PPTX
    intern = state["intern_name"]
    task_date = state.get("task_date") or date.today().strftime("%Y-%m-%d")
    subdir = state.get("subdir", "adhoc")
    topic = state["speech_topic"]

    dot_date = task_date.replace("-", ".")
    filename = general(topic, intern, dot_date, ext="pptx")

    base = Path(config.OUTPUT_DIR) / subdir
    base.mkdir(parents=True, exist_ok=True)
    out_path = base / filename
    prs.save(str(out_path))

    dl_path = Path(config.DOWNLOADS_DIR) / filename
    shutil.copy2(out_path, dl_path)

    # Clean up temp image files
    for slide in structured:
        p = slide.get("image_path")
        if p and Path(p).exists():
            try:
                Path(p).unlink()
            except Exception:
                pass

    # Print non-structured notes to console
    if unstructured:
        print("\n── 非結構化頁面 Notes（需操作者手動處理）" + "─" * 18)
        for i, slide in enumerate(unstructured, 1):
            print(f"\n[非結構 {i}]  {slide['title']}")
            print(slide.get("notes", "（無說明）"))
        print("─" * 50)

    # Write log
    intern_str = _intern_str(intern)
    logger = AgentLogger("speech_ppt_agent", f"Speech PPT: {topic}", intern_str)
    log_path = logger.save(out_path)

    if unstructured:
        with open(log_path, "a", encoding="utf-8") as f:
            f.write("\n\n--- Non-structured Slides (Notes for Operator) ---\n")
            for slide in unstructured:
                f.write(f"\n[{slide['title']}]\n{slide.get('notes', '')}\n")

    print(f"\n  Output: {out_path}")
    print(f"  Log:    {log_path}")

    return {"output_path": str(out_path), "log_path": str(log_path)}


# ── Graph ─────────────────────────────────────────────────────────────────────

def build_graph():
    graph = StateGraph(SpeechPPTState)
    graph.add_node("parse_script",    parse_script)
    graph.add_node("confirm_slides",  confirm_slides)
    graph.add_node("generate_images", generate_images)
    graph.add_node("build_ppt",       build_ppt)

    graph.set_entry_point("parse_script")
    graph.add_edge("parse_script", "confirm_slides")
    graph.add_conditional_edges("confirm_slides", _route_after_confirm, {
        "generate_images": "generate_images",
        END: END,
    })
    graph.add_edge("generate_images", "build_ppt")
    graph.add_edge("build_ppt", END)

    return graph.compile()


# ── Public entry point ────────────────────────────────────────────────────────

def run(
    task_instruction: str,
    speech_topic: str = None,
    intern_name: Union[str, list[str]] = None,
    task_date: str = None,
    subdir: str = "weekly",
    generate_images: bool = True,
    mode: str = "short",
) -> dict:
    """
    Generate a speech PPT from CY's transcript.

    Args:
        task_instruction: CY's verbal description of slide content (transcript).
        speech_topic:     Optional topic label for the filename.  If omitted,
                          derived from the transcript by Claude.
        intern_name:      For file naming.
        task_date:        YYYY-MM-DD; defaults to today.
        subdir:           Output subfolder.
        generate_images:  If False, skip DALL-E (faster, no image cost).
        mode:             Accepted for router compat; not used.

    Returns:
        dict with 'output_path' and 'log_path', or empty dict if cancelled.
    """
    app = build_graph()

    final_state = app.invoke({
        "speech_topic":    speech_topic or "",
        "raw_transcript":  task_instruction or "",
        "intern_name":     intern_name or config.DEFAULT_INTERN_NAME,
        "task_date":       task_date or date.today().strftime("%Y-%m-%d"),
        "subdir":          subdir,
        "generate_images": generate_images,
        "slides_plan":     [],
        "confirmed":       False,
        "output_path":     "",
        "log_path":        "",
    })

    if not final_state.get("output_path"):
        return {}

    return {
        "output_path": final_state["output_path"],
        "log_path":    final_state["log_path"],
    }


# ── CLI ───────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import argparse
    from utils.stt import transcribe

    parser = argparse.ArgumentParser(description="Speech PPT generator")
    parser.add_argument("--topic",     default=None,
                        help="Speech topic label for the filename (optional)")
    parser.add_argument("--audio",     default=None,
                        help="Audio file (.m4a/.mp3/.wav) — transcribed via Whisper")
    parser.add_argument("--text",      default=None,
                        help="Transcript text describing slide content")
    parser.add_argument("--intern",    default=None)
    parser.add_argument("--date",      default=None)
    parser.add_argument("--subdir",    default="weekly",
                        choices=["daily", "weekly", "adhoc"])
    parser.add_argument("--no-images", dest="no_images", action="store_true",
                        help="Skip DALL-E image generation")
    args = parser.parse_args()

    if not args.audio and not args.text:
        parser.error("Provide either --audio or --text.")

    transcript = args.text
    if args.audio:
        print(f"轉錄音檔：{args.audio}")
        transcript = transcribe(args.audio)
        print(f"轉錄完成（{len(transcript)} 字）")

    intern = (
        [n.strip() for n in args.intern.split(",")]
        if args.intern and "," in args.intern
        else args.intern
    )

    result = run(
        task_instruction=transcript,
        speech_topic=args.topic,
        intern_name=intern,
        task_date=args.date,
        subdir=args.subdir,
        generate_images=not args.no_images,
    )

    if not result:
        print("已取消，無輸出。")
    else:
        from utils.cost_tracker import tracker
        tracker.print_task_summary()
