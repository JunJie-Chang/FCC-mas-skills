"""
scripts/build_pptx_cli.py — thin CLI wrapper that builds an FCC-format .pptx
from a JSON slides spec. Optionally invokes OpenAI DALL-E 3 for structured
slide images.

Self-contained: includes layout constants, chrome-template loader, DALL-E
client, and PPTX assembly. The fcc-speech-ppt skill calls this once the
user has confirmed the slides plan.

Usage:
    python3.13 scripts/build_pptx_cli.py --spec /tmp/slides.json
    python3.13 scripts/build_pptx_cli.py --spec -          # stdin

Spec shape:
{
  "speech_topic":   "台中智慧製造",          # required — used in filename
  "task_date":      "2026-05-27",            # optional, default today
  "intern_name":    "Justin",                # optional, default 'Justin'
  "subdir":         "weekly",                # optional, default 'weekly'
  "generate_images": true,                   # optional, default true
  "slides": [
    {"type": "structured",
     "title": "智慧製造的定義",
     "bullets": ["...", "...", "..."]},      # up to 5; extras truncated
    {"type": "unstructured",
     "title": "競爭格局示意圖",
     "notes": "中央放台積電 logo，周邊放上下游廠商…"}
  ]
}

Stdout (JSON):
{
  "output_path":        "/Users/.../output/weekly/2026.05.27_topic_Justin.pptx",
  "unstructured_notes": [
    {"title": "...", "notes": "..."}
  ]
}

Exit 0 on success; 1 on bad spec or failure.
"""
from __future__ import annotations

import argparse
import json
import os
import shutil
import sys
import tempfile
import urllib.request
from datetime import date
from pathlib import Path

# Repo importable
_REPO = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(_REPO))

# Re-exec under the repo venv ($FCC_MAS_HOME/.venv) when launched with a bare
# python3. Must run before any third-party import. See utils/venv_bootstrap.
from utils.venv_bootstrap import activate as _activate_venv  # noqa: E402

_activate_venv(_REPO)

import config  # noqa: E402
from utils.file_naming import general as _fname  # noqa: E402
from utils.spec_io import load_spec, safe_subdir  # noqa: E402

# Load .env so OPENAI_API_KEY is available (DALL-E). override=True keeps
# behavior consistent with the rest of the codebase — see CLAUDE.md.
try:
    from dotenv import load_dotenv
    load_dotenv(_REPO / ".env", override=True)
except Exception:
    pass


# ── Layout constants (inches → EMU: 1in = 914400) ────────────────────────────
_IN = 914400
SLIDE_W = int(10.0 * _IN)
SLIDE_H = int(7.5  * _IN)

TITLE_L, TITLE_T, TITLE_W, TITLE_H, TITLE_PT = (
    int(0.51 * _IN), int(0.40 * _IN), int(8.98 * _IN), int(0.93 * _IN), 28
)
BULLET_L, BULLET_T, BULLET_W, BULLET_H, BULLET_PT = (
    int(0.51 * _IN), int(1.49 * _IN), int(4.70 * _IN), int(4.31 * _IN), 24
)
IMAGE_L, IMAGE_T, IMAGE_W, IMAGE_H = (
    int(5.45 * _IN), int(1.49 * _IN), int(4.20 * _IN), int(4.31 * _IN)
)

# The deep-blue gradient chrome template is committed under assets/ and is
# the single source of truth — no machine-specific reference file.
_CHROME_TEMPLATE = _REPO / "assets" / "ppt_chrome_template.pptx"


# ── Chrome template handling ─────────────────────────────────────────────────

def _ensure_chrome_template() -> bool:
    """Return True if the bundled chrome template asset is present."""
    if _CHROME_TEMPLATE.exists():
        return True
    print(f"[WARN] Missing chrome template asset: {_CHROME_TEMPLATE}", file=sys.stderr)
    return False


def _get_object_layout(prs):
    """Return the 'OBJECT' slide layout; fall back to blank if not found."""
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name == "OBJECT":
            return layout
    return prs.slide_layouts[6]


def _set_bullet_para_props(para) -> None:
    """Apply paragraph formatting matching reference: 150% line spacing, hanging indent, white bullet."""
    from lxml import etree
    from pptx.oxml.ns import qn

    pPr = para._p.get_or_add_pPr()
    pPr.set("marL", "228600")
    pPr.set("indent", "-228600")
    pPr.set("eaLnBrk", "0")
    pPr.set("fontAlgn", "base")
    pPr.set("hangingPunct", "0")

    for tag in ("a:lnSpc", "a:spcBef"):
        old = pPr.find(qn(tag))
        if old is not None:
            pPr.remove(old)

    lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    etree.SubElement(lnSpc, qn("a:spcPct")).set("val", "150000")
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(spcBef, qn("a:spcPts")).set("val", "1000")

    buClr = etree.SubElement(pPr, qn("a:buClr"))
    etree.SubElement(buClr, qn("a:srgbClr")).set("val", "FFFFFF")
    etree.SubElement(pPr, qn("a:buSzPct")).set("val", "100000")
    etree.SubElement(pPr, qn("a:buChar")).set("char", "•")


# ── DALL-E image generation ──────────────────────────────────────────────────

def _generate_dalle_image(title: str, max_retries: int = 3) -> str | None:
    """Generate a DALL-E 3 image for a structured slide title; return local PNG path, or None."""
    try:
        from openai import OpenAI
    except ImportError:
        print("[WARN] openai package not installed; skipping image generation", file=sys.stderr)
        return None

    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        print("[WARN] OPENAI_API_KEY not set; skipping image generation", file=sys.stderr)
        return None

    client = OpenAI(api_key=api_key)
    prompt = (
        f"Photorealistic photograph or realistic illustration representing: '{title}'. "
        "Professional corporate context. Absolutely no text, words, letters, "
        "numbers, signs, labels, or watermarks anywhere in the image. "
        "Suitable as a right-side panel image on a dark navy PowerPoint slide."
    )

    for attempt in range(1, max_retries + 1):
        try:
            resp = client.images.generate(
                model="dall-e-3",
                prompt=prompt,
                size="1024x1024",
                quality="standard",
                n=1,
            )
            image_url = resp.data[0].url
            tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            with urllib.request.urlopen(image_url, timeout=60) as r:
                tmp.write(r.read())
            tmp.close()
            return tmp.name
        except Exception as e:
            if attempt < max_retries:
                print(f"[WARN] retry {attempt}/{max_retries} for '{title}': {e}", file=sys.stderr)
            else:
                print(f"[WARN] image generation failed for '{title}' after {max_retries}", file=sys.stderr)
                return None
    return None


# ── PPTX assembly ────────────────────────────────────────────────────────────

def build_pptx(spec: dict) -> dict:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Pt

    slides_in = spec.get("slides", [])
    structured   = [s for s in slides_in if s.get("type") == "structured"]
    unstructured = [s for s in slides_in if s.get("type") != "structured"]

    # Generate DALL-E images for structured slides (if requested)
    if spec.get("generate_images", True):
        for slide in structured:
            slide["image_path"] = _generate_dalle_image(slide["title"])

    # Load chrome template (or fall back to blank)
    if _ensure_chrome_template():
        prs = Presentation(str(_CHROME_TEMPLATE))
        object_layout = _get_object_layout(prs)
    else:
        prs = Presentation()
        prs.slide_width  = Emu(SLIDE_W)
        prs.slide_height = Emu(SLIDE_H)
        object_layout = prs.slide_layouts[6]

    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    for slide in structured:
        sl = prs.slides.add_slide(object_layout)

        title_box = sl.shapes.add_textbox(Emu(TITLE_L), Emu(TITLE_T), Emu(TITLE_W), Emu(TITLE_H))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = slide["title"]
        run.font.size = Pt(TITLE_PT)
        run.font.bold = False
        run.font.color.rgb = WHITE

        bullet_box = sl.shapes.add_textbox(Emu(BULLET_L), Emu(BULLET_T), Emu(BULLET_W), Emu(BULLET_H))
        btf = bullet_box.text_frame
        btf.word_wrap = True
        for bi, bullet_text in enumerate(slide.get("bullets", [])[:5]):
            bp = btf.paragraphs[0] if bi == 0 else btf.add_paragraph()
            _set_bullet_para_props(bp)
            brun = bp.add_run()
            brun.text = bullet_text
            brun.font.size = Pt(BULLET_PT)
            brun.font.bold = True
            brun.font.color.rgb = WHITE

        img_path = slide.get("image_path")
        if img_path and Path(img_path).exists():
            sl.shapes.add_picture(img_path, Emu(IMAGE_L), Emu(IMAGE_T), Emu(IMAGE_W), Emu(IMAGE_H))

    # Save
    intern   = spec.get("intern_name") or "Justin"
    task_date = spec.get("task_date") or date.today().strftime("%Y-%m-%d")
    subdir   = safe_subdir(spec.get("subdir"), "weekly")
    topic    = spec["speech_topic"]

    dot_date = task_date.replace("-", ".")
    filename = _fname(topic, intern, dot_date, ext="pptx")

    base = Path(config.OUTPUT_DIR) / subdir
    base.mkdir(parents=True, exist_ok=True)
    out_path = base / filename
    prs.save(str(out_path))

    if not getattr(config, "DISABLE_DOWNLOADS_COPY", False):
        dl_path = Path(config.DOWNLOADS_DIR) / filename
        try:
            shutil.copy2(out_path, dl_path)
        except Exception as e:
            print(f"[WARN] Downloads copy failed: {e}", file=sys.stderr)

    # Clean up temp images
    for slide in structured:
        p = slide.get("image_path")
        if p and Path(p).exists():
            try:
                Path(p).unlink()
            except Exception:
                pass

    return {
        "output_path": str(out_path),
        "unstructured_notes": [
            {"title": s["title"], "notes": s.get("notes", "")} for s in unstructured
        ],
    }


def main() -> int:
    ap = argparse.ArgumentParser(description="Build an FCC-format .pptx from JSON spec.")
    ap.add_argument("--spec", required=True, help="Path to JSON spec, or '-' for stdin.")
    args = ap.parse_args()

    try:
        spec = load_spec(args.spec)
    except json.JSONDecodeError as e:
        print(f"ERROR: invalid JSON spec — {e}", file=sys.stderr)
        return 1

    if "speech_topic" not in spec:
        print("ERROR: spec missing required field 'speech_topic'", file=sys.stderr)
        return 1
    if not spec.get("slides"):
        print("ERROR: spec.slides is empty", file=sys.stderr)
        return 1

    try:
        result = build_pptx(spec)
    except Exception as e:
        print(f"ERROR: PPTX build failed — {e}", file=sys.stderr)
        return 1

    print(json.dumps(result, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    sys.exit(main())
